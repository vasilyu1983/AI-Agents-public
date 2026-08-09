#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any


def _require_python_pptx():
    try:
        from pptx import Presentation  # type: ignore

        return Presentation
    except ImportError as exc:
        raise RuntimeError(
            "Missing dependency: python-pptx. Install with: pip install python-pptx"
        ) from exc


def _slide_title(slide: Any, index: int) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "text", "").strip():
        return title_shape.text.strip()
    return f"Slide {index}"


def extract_notes(pptx_path: Path) -> list[dict[str, Any]]:
    Presentation = _require_python_pptx()
    prs = Presentation(str(pptx_path))

    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append(
            {
                "index": index,
                "title": _slide_title(slide, index),
                "has_notes": bool(notes),
                "notes": notes,
            }
        )
    return slides


def to_markdown(slides: list[dict[str, Any]]) -> str:
    lines = ["# Speaker Notes", ""]
    for slide in slides:
        lines.append(f"## Slide {slide['index']}: {slide['title']}")
        lines.append("")
        lines.append(slide["notes"] or "_No notes._")
        lines.append("")
    return "\n".join(lines)


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description="Extract PPTX speaker notes to JSON or Markdown.")
    parser.add_argument("pptx", type=Path, help="Path to a .pptx file")
    parser.add_argument(
        "--format",
        choices=("json", "markdown"),
        default="json",
        help="Output format (default: json)",
    )
    parser.add_argument("--out", type=Path, help="Output path (defaults to stdout)")
    args = parser.parse_args(argv)

    if not args.pptx.exists():
        print(f"File not found: {args.pptx}", file=sys.stderr)
        return 2

    if args.pptx.suffix.lower() != ".pptx":
        print("Expected a .pptx file.", file=sys.stderr)
        return 2

    try:
        slides = extract_notes(args.pptx)
    except Exception as exc:
        print(str(exc), file=sys.stderr)
        return 2

    if args.format == "markdown":
        output = to_markdown(slides)
    else:
        output = json.dumps(
            {
                "path": str(args.pptx),
                "slides": slides,
            },
            indent=2,
            ensure_ascii=False,
        )

    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(output, encoding="utf-8")
    else:
        print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
