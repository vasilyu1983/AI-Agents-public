#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any


def _require_python_pptx():
    try:
        from pptx import Presentation  # type: ignore

        return Presentation
    except ImportError as exc:
        raise RuntimeError(
            "Missing dependency: python-pptx. Install with: pip install python-pptx"
        ) from exc


def _shape_kind(shape: Any) -> str:
    shape_type = getattr(shape, "shape_type", None)
    return getattr(shape_type, "name", str(shape_type))


def _slide_title(slide: Any, index: int) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and getattr(title_shape, "text", "").strip():
        return title_shape.text.strip()
    return f"Slide {index}"


def inventory_pptx(pptx_path: Path) -> dict[str, Any]:
    Presentation = _require_python_pptx()
    prs = Presentation(str(pptx_path))

    masters = []
    slide_masters = list(getattr(prs, "slide_masters", []))
    if not slide_masters and getattr(prs, "slide_master", None) is not None:
        slide_masters = [prs.slide_master]

    for master_index, master in enumerate(slide_masters):
        layouts = []
        for layout_index, layout in enumerate(master.slide_layouts):
            placeholders = []
            for placeholder in layout.placeholders:
                placeholders.append(
                    {
                        "idx": placeholder.placeholder_format.idx,
                        "name": placeholder.name,
                        "type": str(placeholder.placeholder_format.type),
                    }
                )
            layouts.append(
                {
                    "index": layout_index,
                    "name": layout.name,
                    "placeholders": placeholders,
                }
            )
        masters.append({"index": master_index, "layouts": layouts})

    slides = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        shapes = []
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        for shape in slide.shapes:
            entry = {
                "name": getattr(shape, "name", None),
                "type": _shape_kind(shape),
                "has_text": bool(getattr(shape, "has_text_frame", False)),
                "has_table": bool(getattr(shape, "has_table", False)),
                "has_chart": bool(getattr(shape, "has_chart", False)),
            }
            if getattr(shape, "is_placeholder", False):
                placeholder_format = shape.placeholder_format
                entry["placeholder_idx"] = getattr(placeholder_format, "idx", None)
                entry["placeholder_type"] = str(getattr(placeholder_format, "type", ""))
            shapes.append(entry)

        slides.append(
            {
                "index": slide_index,
                "title": _slide_title(slide, slide_index),
                "shape_count": len(slide.shapes),
                "has_notes": bool(notes),
                "shapes": shapes,
            }
        )

    return {
        "path": str(pptx_path),
        "slide_count": len(prs.slides),
        "master_count": len(masters),
        "masters": masters,
        "slides": slides,
    }


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(
        description="Inventory PPTX masters, layouts, placeholders, slides, and notes presence."
    )
    parser.add_argument("pptx", type=Path, help="Path to a .pptx file")
    parser.add_argument("--json", action="store_true", help="Emit JSON to stdout")
    args = parser.parse_args(argv)

    if not args.pptx.exists():
        print(f"File not found: {args.pptx}", file=sys.stderr)
        return 2

    if args.pptx.suffix.lower() != ".pptx":
        print("Expected a .pptx file.", file=sys.stderr)
        return 2

    try:
        payload = inventory_pptx(args.pptx)
    except Exception as exc:
        print(str(exc), file=sys.stderr)
        return 2

    if args.json:
        print(json.dumps(payload, indent=2, ensure_ascii=False))
        return 0

    print(f"File: {payload['path']}")
    print(f"Slides: {payload['slide_count']}")
    print(f"Slide masters: {payload['master_count']}")
    print("")
    print("Layouts:")
    for master in payload["masters"]:
        for layout in master["layouts"]:
            print(f"- {layout['name']} ({len(layout['placeholders'])} placeholders)")

    print("")
    print("Slides:")
    for slide in payload["slides"]:
        suffix = " with notes" if slide["has_notes"] else ""
        print(f"- {slide['index']}: {slide['title']} ({slide['shape_count']} shapes){suffix}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
